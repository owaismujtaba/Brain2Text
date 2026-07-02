"""Detailed NJEM deck — data, architecture, losses (incl. decoder-in-the-loop),
training objective, evaluation, and current status. Produces NJEM_detailed.pptx."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ── Modern light palette ────────────────────────────────
INK    = RGBColor(0x0F, 0x17, 0x2A)
SLATE  = RGBColor(0x33, 0x41, 0x55)
MUTED  = RGBColor(0x64, 0x74, 0x8B)
BORDER = RGBColor(0xE2, 0xE8, 0xF0)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
BG_A   = RGBColor(0xFB, 0xFC, 0xFF)
BG_B   = RGBColor(0xEA, 0xF0, 0xFD)
CARD_A = RGBColor(0xFF, 0xFF, 0xFF)
CARD_B = RGBColor(0xF3, 0xF6, 0xFD)
INK_A  = RGBColor(0x1E, 0x1B, 0x4B)   # deep indigo (dark panels)
INK_B  = RGBColor(0x31, 0x2E, 0x81)
ACCENTS = [
    (RGBColor(0x63,0x66,0xF1), RGBColor(0x81,0x8C,0xF8), RGBColor(0x4F,0x46,0xE5)),  # indigo
    (RGBColor(0x8B,0x5C,0xF6), RGBColor(0xA7,0x8B,0xFA), RGBColor(0x7C,0x3A,0xED)),  # violet
    (RGBColor(0x06,0xB6,0xD4), RGBColor(0x22,0xD3,0xEE), RGBColor(0x08,0x91,0xB2)),  # cyan
    (RGBColor(0x14,0xB8,0xA6), RGBColor(0x2D,0xD4,0xBF), RGBColor(0x0D,0x94,0x88)),  # teal
    (RGBColor(0xF5,0x9E,0x0B), RGBColor(0xFB,0xBF,0x24), RGBColor(0xD9,0x77,0x06)),  # amber
    (RGBColor(0xEC,0x48,0x99), RGBColor(0xF4,0x72,0xB6), RGBColor(0xDB,0x27,0x77)),  # pink
]

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = 13.333, 7.5

def grad(fill, c1, c2, angle=90):
    fill.gradient()
    s = fill.gradient_stops
    s[0].position = 0.0; s[0].color.rgb = c1
    s[1].position = 1.0; s[1].color.rgb = c2
    try: fill.gradient_angle = angle
    except Exception: pass

def new_slide(dark=False):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    if dark:
        grad(sl.background.fill, INK_A, INK_B, angle=115)
    else:
        grad(sl.background.fill, BG_A, BG_B, angle=115)
    return sl

def txt(sl, x, y, w, h, lines, size, color, bold=False, italic=False,
        align=PP_ALIGN.LEFT, font="Calibri", anchor=MSO_ANCHOR.TOP, spacing=1.0, cspace=None):
    tb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
    if isinstance(lines, str): lines = [lines]
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = spacing
        r = p.add_run(); r.text = ln
        f = r.font; f.size = Pt(size); f.bold = bold; f.italic = italic
        f.name = font; f.color.rgb = color
        if cspace is not None:
            r._r.get_or_add_rPr().set("spc", str(cspace))
    return tb

def soft_shadow(sp):
    spPr = sp._element.spPr
    eff = spPr.makeelement(qn('a:effectLst'), {})
    sh = spPr.makeelement(qn('a:outerShdw'),
                          {'blurRad':'90000','dist':'40000','dir':'5400000','rotWithShape':'0'})
    clr = spPr.makeelement(qn('a:srgbClr'), {'val':'9FB0CC'})
    clr.append(spPr.makeelement(qn('a:alpha'), {'val':'38000'}))
    sh.append(clr); eff.append(sh); spPr.append(eff)

def card(sl, x, y, w, h, c1=CARD_A, c2=CARD_B, angle=90, radius=0.10, border=BORDER, shadow=True):
    sp = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    grad(sp.fill, c1, c2, angle)
    if border is not None: sp.line.color.rgb = border; sp.line.width = Pt(1)
    else: sp.line.fill.background()
    try: sp.adjustments[0] = radius
    except Exception: pass
    sp.shadow.inherit = False
    if shadow: soft_shadow(sp)
    return sp

def circle(sl, x, y, d, c1, c2):
    sp = sl.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    grad(sp.fill, c1, c2, angle=45); sp.line.fill.background(); sp.shadow.inherit = False
    soft_shadow(sp); return sp

def header(sl, title, subtitle, accent=0):
    txt(sl, 0.6, 0.4, 12.1, 0.7, title, 32, INK, bold=True, font="Cambria")
    txt(sl, 0.62, 1.1, 12.1, 0.42, subtitle, 15, ACCENTS[accent][2], italic=True)

def footer(sl, text):
    txt(sl, 0.6, 7.02, 12.1, 0.35, text, 10.5, MUTED, align=PP_ALIGN.CENTER)

def arrow(sl, x, y, w, h, ch, color, size=22):
    txt(sl, x, y, w, h, ch, size, color, bold=True,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# ════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ════════════════════════════════════════════════════════
s = new_slide(dark=True)
txt(s, 0.9, 1.35, 11.5, 0.4, "NJEM  ·  NEURAL JOINT EMBEDDING MODEL", 15,
    RGBColor(0xA7,0x8B,0xFA), bold=True, cspace=280)
txt(s, 0.88, 1.95, 11.6, 1.5, ["Decoding Speech from", "Brain Activity"], 52, WHITE,
    bold=True, font="Cambria", spacing=0.98)
txt(s, 0.9, 3.95, 11.0, 0.5,
    "Mapping intracortical neural signals to Whisper audio embeddings, then to text",
    18, RGBColor(0xCB,0xD5,0xF5), italic=True)

stats = [
    ("12.3M", "ConvBiGRU parameters"),
    ("1500 x 384", "Whisper embedding target"),
    ("8,072 / 1,425", "train / val trials"),
]
cW, cGap = 3.62, 0.35
x0 = 0.9
sy = 4.95
for i, (big, lab) in enumerate(stats):
    a1, a2, adk = ACCENTS[i]
    x = x0 + i*(cW+cGap)
    cd = card(s, x, sy, cW, 1.35, c1=RGBColor(0x2A,0x27,0x63), c2=RGBColor(0x22,0x1F,0x52),
              border=None);
    txt(s, x+0.32, sy+0.22, cW-0.5, 0.55, big, 27, WHITE, bold=True, font="Cambria")
    txt(s, x+0.34, sy+0.86, cW-0.5, 0.4, lab, 13, RGBColor(0xA7,0xB0,0xE0))
txt(s, 0.9, 6.72, 11.5, 0.4, "t15 Brain-to-Text dataset   |   frozen Whisper tiny.en decoder",
    13, RGBColor(0x8A,0x93,0xC8))

# ════════════════════════════════════════════════════════
# SLIDE 2 — Pipeline overview
# ════════════════════════════════════════════════════════
s = new_slide()
header(s, "The Pipeline — Four Stages", "From raw cortical recordings to a text transcript", 0)
stages = [
    ("Neural Activity", "(T x 512)", ["Intracortical t15", "recordings, 20 ms bins"]),
    ("ConvBiGRU Model", "12.3M params", ["Conv + 8-layer BiGRU", "+ cross-attention"]),
    ("Whisper Embeddings", "(1500 x 384)", ["Predicted encoder", "feature sequence"]),
    ("Text Transcript", "frozen decoder", ["Whisper tiny.en decodes", "embeddings -> words"]),
]
cardW, cardH, gap = 2.75, 2.35, 0.42
x0 = (SW - (4*cardW + 3*gap)) / 2
yC = 2.15
for i, (t, sp_, d) in enumerate(stages):
    a1, a2, adk = ACCENTS[i]
    x = x0 + i*(cardW+gap)
    card(s, x, yC, cardW, cardH)
    cd = 0.66
    circle(s, x+cardW/2-cd/2, yC+0.26, cd, a1, a2)
    txt(s, x+cardW/2-cd/2, yC+0.26, cd, cd, str(i+1), 22, WHITE, bold=True,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font="Cambria")
    txt(s, x+0.1, yC+1.05, cardW-0.2, 0.55, t, 15, INK, bold=True, align=PP_ALIGN.CENTER, font="Cambria")
    txt(s, x, yC+1.58, cardW, 0.3, sp_, 12, adk, bold=True, align=PP_ALIGN.CENTER, font="Courier New")
    txt(s, x+0.12, yC+1.92, cardW-0.24, 0.7, d, 11, MUTED, align=PP_ALIGN.CENTER, spacing=0.95)
    if i < 3:
        arrow(s, x+cardW-0.02, yC+0.9, gap+0.06, 0.6, "->", adk, 22)

facts = [
    ("Why embeddings?", ["Whisper's encoder space is", "already speech-structured"]),
    ("Frozen decoder", ["Reuse Whisper tiny.en to turn", "embeddings into words"]),
    ("Direct WER signal", ["Decoder-in-the-loop trains", "for decodability, not just MSE"]),
]
fW, fGap = 3.9, 0.3
fx0 = (SW - (3*fW + 2*fGap)) / 2
fy = 5.05
for i, (k, v) in enumerate(facts):
    a1, a2, adk = ACCENTS[i]
    x = fx0 + i*(fW+fGap)
    card(s, x, fy, fW, 1.5)
    circle(s, x+0.28, fy+0.3, 0.34, a1, a2)
    txt(s, x+0.78, fy+0.24, fW-0.95, 0.4, k, 14, INK, bold=True, font="Cambria")
    txt(s, x+0.3, fy+0.72, fW-0.55, 0.7, v, 12, MUTED, spacing=0.95)
footer(s, "NJEM — Neural Joint Embedding Model")

# ════════════════════════════════════════════════════════
# SLIDE 3 — Data & preprocessing
# ════════════════════════════════════════════════════════
s = new_slide()
header(s, "Data & Preprocessing", "Pairing neural recordings with Whisper encoder targets per (session, trial)", 2)

left = [
    ("Neural input  x", "(T, 512) float32",
     "512 features per 20 ms bin, z-scored intracortical activity from the t15 array."),
    ("Target embedding", "(vf, 384) float32",
     "Whisper tiny.en encoder output, kept only for valid content frames."),
    ("Transcription", "str",
     "Ground-truth sentence — drives both WER eval and the decoder-in-the-loop loss."),
]
lx, lw = 0.62, 6.35
ly = 1.75; lh = 1.5; lgp = 0.16
for i, (t, shp, body) in enumerate(left):
    a1, a2, adk = ACCENTS[i]
    y = ly + i*(lh+lgp)
    card(s, lx, y, lw, lh)
    circle(s, lx+0.3, y+0.32, 0.4, a1, a2)
    txt(s, lx+0.92, y+0.24, lw-1.1, 0.35, t, 15, INK, bold=True, font="Cambria")
    txt(s, lx+0.92, y+0.62, lw-1.1, 0.3, shp, 11, adk, bold=True, font="Courier New")
    txt(s, lx+0.32, y+1.0, lw-0.6, 0.42, body, 11.5, MUTED, spacing=0.95)

# right column: alignment + collate + sizes
rx, rw = 7.35, 5.35
card(s, rx, 1.75, rw, 1.62)
txt(s, rx+0.3, 1.9, rw-0.6, 0.35, "Frame alignment (length masking)", 15, INK, bold=True, font="Cambria")
txt(s, rx+0.3, 2.32, rw-0.6, 0.95,
    ["vf = clamp( ceil(audio_length / 320), 1, 1500 )",
     "320 samples = 20 ms @ 16 kHz  ->  one encoder frame",
     "Only these vf frames count in every loss term."],
    11, ACCENTS[2][2], font="Courier New", spacing=1.15)

card(s, rx, 3.55, rw, 1.5)
txt(s, rx+0.3, 3.7, rw-0.6, 0.35, "Collation (per batch)", 15, INK, bold=True, font="Cambria")
txt(s, rx+0.3, 4.12, rw-0.6, 0.9,
    ["neural  -> pad to max T           (B, T, 512)",
     "target  -> zero-pad to 1500       (B, 1500, 384)",
     "mask    -> bool, True on vf       (B, 1500)"],
    10.5, MUTED, font="Courier New", spacing=1.15)

# dataset size chips
chips = [("8,072", "train trials", 0), ("1,425", "val trials", 3)]
cw2 = 2.55
for i, (big, lab, ai) in enumerate(chips):
    a1, a2, adk = ACCENTS[ai]
    x = rx + i*(cw2+0.25)
    card(s, x, 5.25, cw2, 1.05, c1=CARD_A, c2=CARD_B)
    txt(s, x+0.28, 5.42, cw2-0.4, 0.5, big, 24, adk, bold=True, font="Cambria")
    txt(s, x+0.3, 5.98, cw2-0.4, 0.3, lab, 12, MUTED)
footer(s, "src/dataset.py  —  NeuralToEmbeddingDataset  ·  valid_frames()  ·  collate()")

# ════════════════════════════════════════════════════════
# SLIDE 4 — Architecture forward pass
# ════════════════════════════════════════════════════════
s = new_slide()
header(s, "Model Architecture — ConvBiGRU", "Forward pass: neural sequence -> aligned 384-dim embedding per frame", 0)

blocks = [
    ("Input — neural features", "x : (B, T, 512)   z-scored, 20 ms bins"),
    ("Conv front-end", "2 x Conv1d(k=5, s=2) + GELU + Dropout  ->  (B, T/4, 256)"),
    ("BiGRU (packed sequences)", "8 layers · hidden 256 · bidirectional  ->  (B, T', 512)"),
    ("Cross-attention", "1500 learnable queries · 4 heads · key-pad mask  ->  (B, 1500, 512)"),
    ("LayerNorm + MLP head", "Linear 512->512, GELU, Linear 512->384  ->  (B, 1500, 384)"),
    ("Frozen Whisper decoder", "predicted embeddings  ->  text tokens   (inference)"),
]
lx, lw = 0.62, 7.35
by = 1.72; bh = 0.72; bgap = 0.10
for i, (t, sub) in enumerate(blocks):
    a1, a2, adk = ACCENTS[i]
    y = by + i*(bh+bgap)
    card(s, lx, y, lw, bh)
    cd = 0.44
    circle(s, lx+0.22, y+bh/2-cd/2, cd, a1, a2)
    txt(s, lx+0.22, y+bh/2-cd/2, cd, cd, str(i+1), 15, WHITE, bold=True,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font="Cambria")
    txt(s, lx+0.85, y+0.11, lw-1.0, 0.32, t, 14, INK, bold=True, font="Cambria")
    txt(s, lx+0.85, y+0.43, lw-1.0, 0.28, sub, 10.5, MUTED, font="Courier New")

# right notes
rx, rw = 8.28, 4.45
notes = [
    ("Why cross-attention?", ["Neural length T/4 varies per trial;", "1500 learnable queries resample it", "to Whisper's fixed frame grid."]),
    ("Padding handled twice", ["packed GRU + attention key-mask", "so padded frames never leak into", "the output embeddings."]),
    ("Regression head", ["per-frame MLP predicts the raw", "384-dim encoder vector (no", "softmax) — it is a regressor."]),
]
ny = 1.72; nh = 1.62; ngp = 0.14
for i, (t, body) in enumerate(notes):
    a1, a2, adk = ACCENTS[i]
    y = ny + i*(nh+ngp)
    card(s, rx, y, rw, nh)
    txt(s, rx+0.3, y+0.2, rw-0.55, 0.35, t, 13.5, INK, bold=True, font="Cambria")
    txt(s, rx+0.3, y+0.6, rw-0.55, 0.95, body, 11, MUTED, spacing=1.0)
footer(s, "src/model.py  —  ConvBiGRU  ·  in_dim 512  ·  emb_dim 384  ·  n_ctx 1500")

# ════════════════════════════════════════════════════════
# SLIDE 5 — Embedding loss
# ════════════════════════════════════════════════════════
s = new_slide()
header(s, "Loss 1 — Masked Embedding Regression", "Match predicted embeddings to Whisper targets, over valid frames only", 3)

# formula panel
card(s, 0.62, 1.72, 7.2, 2.5, c1=INK_A, c2=INK_B, border=None)
txt(s, 0.95, 1.95, 6.6, 0.4, "masked_embedding_loss", 15, RGBColor(0xA7,0x8B,0xFA), bold=True, font="Courier New")
txt(s, 0.95, 2.5, 6.6, 1.6,
    ["L_emb =  L1_w · SmoothL1(pred, tgt)",
     "              +  cos_w · ( 1 - cos(pred, tgt) )",
     "",
     "averaged over valid frames  (mask = True)",
     "L1_w = cos_w = 1.0"],
    14, WHITE, font="Courier New", spacing=1.2)

# two rationale cards
rat = [
    ("SmoothL1 term", ACCENTS[0][2],
     ["Robust magnitude match — less", "sensitive to outlier dimensions", "than plain MSE."]),
    ("Cosine term", ACCENTS[2][2],
     ["Aligns direction of each 384-d", "vector. The Whisper decoder is", "highly sensitive to direction."]),
]
rx, rw = 8.28, 4.45
for i, (t, col, body) in enumerate(rat):
    y = 1.72 + i*1.28
    card(s, rx, y, rw, 1.14)
    txt(s, rx+0.3, y+0.16, rw-0.55, 0.35, t, 14, INK, bold=True, font="Cambria")
    txt(s, rx+0.3, y+0.55, rw-0.55, 0.55, body, 11, MUTED, spacing=1.0)

# the gap callout
card(s, 0.62, 4.5, 12.1, 1.9, c1=CARD_A, c2=CARD_B)
circle(s, 0.95, 4.78, 0.42, ACCENTS[4][0], ACCENTS[4][1])
txt(s, 1.55, 4.72, 11.0, 0.4, "The problem this alone cannot solve", 16, INK, bold=True, font="Cambria")
txt(s, 1.0, 5.3, 11.5, 1.0,
    ["A low regression loss does NOT guarantee a low WER. The model can minimise SmoothL1/cosine by",
     "regressing toward the average embedding (mode collapse): cosine plateaus near ~0.72 while the decoder",
     "needs ~0.95+ to recover the right words. Symptom seen in training: fluent but wrong sentences, WER ~ 1.0."],
    12.5, SLATE, spacing=1.1)
footer(s, "src/model.py  —  masked_embedding_loss(pred, target, mask)")

# ════════════════════════════════════════════════════════
# SLIDE 6 — Decoder-in-the-loop loss
# ════════════════════════════════════════════════════════
s = new_slide()
header(s, "Loss 2 — Decoder-in-the-Loop (the key idea)", "Optimise decodability directly by back-propagating through the frozen decoder", 1)

# flow chain
nodes = [
    ("Predicted\nembeddings", "(B, T, 384)", ACCENTS[0]),
    ("FROZEN Whisper\ndecoder", "teacher-forced", ACCENTS[1]),
    ("Token logits", "(B, L, vocab)", ACCENTS[2]),
    ("Cross-entropy\nvs GT tokens", "L_CE", ACCENTS[4]),
]
nW, nGap = 2.72, 0.55
nx0 = (SW - (4*nW + 3*nGap)) / 2
ny = 1.78
for i, (t, sub, acc) in enumerate(nodes):
    a1, a2, adk = acc
    x = nx0 + i*(nW+nGap)
    card(s, x, ny, nW, 1.35)
    txt(s, x+0.2, ny+0.24, nW-0.4, 0.6, t.split("\n"), 14, INK, bold=True,
        align=PP_ALIGN.CENTER, font="Cambria", spacing=0.95)
    txt(s, x+0.2, ny+0.95, nW-0.4, 0.3, sub, 11, adk, bold=True,
        align=PP_ALIGN.CENTER, font="Courier New")
    if i < 3:
        arrow(s, x+nW, ny+0.4, nGap, 0.55, "->", MUTED, 22)

# mechanics (left) + why (right)
lx, lw = 0.62, 7.2
card(s, lx, 3.5, lw, 2.9, c1=INK_A, c2=INK_B, border=None)
txt(s, lx+0.32, 3.68, lw-0.6, 0.35, "How it is computed", 15, RGBColor(0xA7,0x8B,0xFA), bold=True, font="Cambria")
txt(s, lx+0.32, 4.12, lw-0.6, 2.2,
    ["feats = pred * mask_vf          # zero out padding frames",
     "tokens = SOT + encode(' '+text) + EOT",
     "inp, tgt = tokens[:-1], tokens[1:]   # shift by one",
     "tgt[prompt] = -100               # ignore prompt positions",
     "logits = whisper.decoder(inp, feats)",
     "L_CE = cross_entropy(logits, tgt, ignore_index=-100)"],
    11.5, WHITE, font="Courier New", spacing=1.28)

rx, rw = 8.1, 4.62
whys = [
    ("Decoder stays frozen", ["No decoder weights update — only", "the gradient flows back into the", "predicted embeddings."]),
    ("Trains for decodability", ["Pushes embeddings into the region", "Whisper actually reads as the right", "words — closes the MSE<->WER gap."]),
]
for i, (t, body) in enumerate(whys):
    a1, a2, adk = ACCENTS[i+2]
    y = 3.5 + i*1.5
    card(s, rx, y, rw, 1.36)
    circle(s, rx+0.28, y+0.24, 0.34, a1, a2)
    txt(s, rx+0.76, y+0.18, rw-0.95, 0.35, t, 13.5, INK, bold=True, font="Cambria")
    txt(s, rx+0.3, y+0.6, rw-0.55, 0.7, body, 11, MUTED, spacing=1.0)
footer(s, "src/decoder_loss.py  —  WhisperDecoderLoss  ·  weight = dec_loss_weight = 0.8")

# ════════════════════════════════════════════════════════
# SLIDE 7 — Total objective & training config
# ════════════════════════════════════════════════════════
s = new_slide()
header(s, "Training Objective & Configuration", "One combined loss, AdamW with cosine annealing, robust checkpointing", 0)

# total loss banner
card(s, 0.62, 1.72, 12.1, 1.15, c1=INK_A, c2=INK_B, border=None)
txt(s, 0.95, 1.9, 11.4, 0.4, "Total loss", 14, RGBColor(0xA7,0x8B,0xFA), bold=True, font="Cambria")
txt(s, 0.95, 2.3, 11.4, 0.5,
    "L  =  L_emb  +  λ · L_CE          λ = dec_loss_weight = 0.8",
    17, WHITE, font="Courier New")

# config grid
cfg = [
    ("epochs", "50"), ("batch_size", "12"), ("optimizer", "AdamW"),
    ("lr", "3e-4"), ("weight_decay", "1e-4"), ("scheduler", "cosine"),
    ("gru_layers", "8"), ("hidden", "256"), ("conv_channels", "256"),
    ("dropout", "0.1"), ("grad clip", "1.0"), ("dec_model", "tiny.en"),
]
gcols, grows = 3, 4
gW, gH, gxg, gyg = 3.9, 0.62, 0.2, 0.16
gx0 = 0.62; gy0 = 3.12
for idx, (k, v) in enumerate(cfg):
    r, c = divmod(idx, gcols)
    x = gx0 + c*(gW+gxg); y = gy0 + r*(gH+gyg)
    card(s, x, y, gW, gH, radius=0.16)
    txt(s, x+0.28, y+0.13, gW-1.4, 0.36, k, 12.5, SLATE, font="Courier New", anchor=MSO_ANCHOR.MIDDLE)
    txt(s, x+gW-1.55, y+0.13, 1.3, 0.36, v, 13.5, ACCENTS[0][2], bold=True,
        align=PP_ALIGN.RIGHT, font="Courier New", anchor=MSO_ANCHOR.MIDDLE)

# checkpointing strip
card(s, 0.62, 6.35, 12.1, 0.85)
txt(s, 0.9, 6.5, 2.4, 0.5, "Checkpoints", 13.5, INK, bold=True, font="Cambria", anchor=MSO_ANCHOR.MIDDLE)
ck = [("best.pt", "lowest val loss", 0), ("best_wer.pt", "lowest WER", 3), ("last.pt", "every epoch (resume)", 1)]
cx = 3.2
for i, (n, d, ai) in enumerate(ck):
    x = cx + i*3.15
    txt(s, x, 6.46, 3.0, 0.3, n, 12, ACCENTS[ai][2], bold=True, font="Courier New")
    txt(s, x, 6.76, 3.0, 0.3, d, 10.5, MUTED)
footer(s, "config.yaml  ·  src/train_neural2emb.py  —  resume prefers last.pt, then best.pt")

# ════════════════════════════════════════════════════════
# SLIDE 8 — Evaluation & decoding
# ════════════════════════════════════════════════════════
s = new_slide()
header(s, "Evaluation & Decoding", "Turn predicted embeddings into words and score against ground truth", 2)

steps = [
    ("Skip the encoder", "whisper.decode( feats )",
     "Feed the predicted (1500, 384) tensor straight in as if it were encoder output."),
    ("Beam search", "beam_size = 5",
     "Beam decoding with language='en', without_timestamps for stable transcripts."),
    ("Normalize + WER", "EnglishTextNormalizer",
     "Levenshtein word error rate on normalized text; also track exact-match %."),
    ("Write transcripts", "outputs/*.csv",
     "decode_dataset -> rows of (idx, wer, truth, pred) for inspection."),
]
sx, sw = 0.62, 6.1
sy = 1.75; shh = 1.22; sgp = 0.16
for i, (t, code, body) in enumerate(steps):
    a1, a2, adk = ACCENTS[i]
    y = sy + i*(shh+sgp)
    card(s, sx, y, sw, shh)
    circle(s, sx+0.3, y+0.3, 0.4, a1, a2)
    txt(s, sx+0.3, y+0.3, 0.4, 0.4, str(i+1), 15, WHITE, bold=True,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font="Cambria")
    txt(s, sx+0.92, y+0.16, sw-1.1, 0.34, t, 14, INK, bold=True, font="Cambria")
    txt(s, sx+0.92, y+0.52, sw-1.1, 0.3, code, 10.5, adk, bold=True, font="Courier New")
    txt(s, sx+0.32, y+0.84, sw-0.6, 0.34, body, 11, MUTED, spacing=0.95)

# right: cadence + entry points
rx, rw = 7.05, 5.65
card(s, rx, 1.75, rw, 1.5)
txt(s, rx+0.3, 1.9, rw-0.55, 0.35, "Eval cadence during training", 15, INK, bold=True, font="Cambria")
txt(s, rx+0.3, 2.32, rw-0.55, 0.85,
    ["wer_every = 5     -> WER eval every 5 epochs",
     "wer_trials = 30   -> val trials scored per eval",
     "best_wer.pt saved whenever WER improves"],
    11, MUTED, font="Courier New", spacing=1.2)

card(s, rx, 3.4, rw, 1.5)
txt(s, rx+0.3, 3.55, rw-0.55, 0.35, "Standalone decode stage", 15, INK, bold=True, font="Cambria")
txt(s, rx+0.3, 3.97, rw-0.55, 0.85,
    ["python main.py decode",
     "loads a checkpoint, rebuilds ConvBiGRU,",
     "decodes a split, writes the CSV of results."],
    11, ACCENTS[3][2], font="Courier New", spacing=1.2)

card(s, rx, 5.05, rw, 1.35, c1=INK_A, c2=INK_B, border=None)
txt(s, rx+0.3, 5.2, rw-0.55, 0.35, "Why decode from embeddings works", 13.5,
    RGBColor(0xA7,0x8B,0xFA), bold=True, font="Cambria")
txt(s, rx+0.3, 5.6, rw-0.55, 0.7,
    ["Whisper's decoder cross-attends to encoder", "features; supplying our predicted features",
     "bypasses audio entirely."], 11, WHITE, spacing=1.05)
footer(s, "src/decode.py  —  decode_dataset()  ·  decode()")

# ════════════════════════════════════════════════════════
# SLIDE 9 — Current status & results
# ════════════════════════════════════════════════════════
s = new_slide()
header(s, "Current Status & Results", "Decoder-in-the-loop is producing fluent output; content accuracy is the next milestone", 5)

# metric chips
mets = [("epoch 7", "of 50", 0), ("3.47", "val loss (down)", 3), ("1.02", "WER (val)", 4), ("0%", "exact match", 1)]
mW = 2.92
for i, (big, lab, ai) in enumerate(mets):
    a1, a2, adk = ACCENTS[ai]
    x = 0.62 + i*(mW+0.18)
    card(s, x, 1.72, mW, 1.1)
    txt(s, x+0.28, 1.86, mW-0.5, 0.5, big, 24, adk, bold=True, font="Cambria")
    txt(s, x+0.3, 2.42, mW-0.5, 0.3, lab, 12, MUTED)

# sample decodings
card(s, 0.62, 3.05, 7.55, 3.15)
txt(s, 0.9, 3.2, 7.0, 0.35, "Sample decodings (epoch 5 eval)", 14, INK, bold=True, font="Cambria")
rows = [
    ("You can see the code at this point as well.", "You know, you know, you know, you know."),
    ("How does it keep the cost down?", "I don't know how to do it."),
    ("Not too controversial.", "That's it."),
]
ry = 3.68
for tr, pr in rows:
    txt(s, 0.95, ry, 0.7, 0.3, "truth", 10, ACCENTS[3][2], bold=True, font="Courier New")
    txt(s, 1.7, ry, 6.3, 0.3, tr, 11.5, SLATE)
    txt(s, 0.95, ry+0.34, 0.7, 0.3, "pred", 10, ACCENTS[4][2], bold=True, font="Courier New")
    txt(s, 1.7, ry+0.34, 6.3, 0.3, pr, 11.5, MUTED, italic=True)
    ry += 0.82

# interpretation + next steps
rx, rw = 8.3, 4.42
card(s, rx, 3.05, rw, 1.5, c1=CARD_A, c2=CARD_B)
txt(s, rx+0.28, 3.2, rw-0.5, 0.35, "What this tells us", 13.5, INK, bold=True, font="Cambria")
txt(s, rx+0.28, 3.6, rw-0.5, 0.9,
    ["Output is now fluent English (decoder loss",
     "working) but content-wrong — classic early",
     "collapse. dec-CE is falling steadily each epoch."],
    11, MUTED, spacing=1.02)

card(s, rx, 4.7, rw, 1.5, c1=INK_A, c2=INK_B, border=None)
txt(s, rx+0.28, 4.85, rw-0.5, 0.35, "Next steps", 13.5, RGBColor(0xA7,0x8B,0xFA), bold=True, font="Cambria")
txt(s, rx+0.28, 5.25, rw-0.5, 0.9,
    ["- run full 50 epochs (single GPU process)",
     "- watch dec-CE + cosine toward ~0.95",
     "- resume from last.pt (epoch 8)"],
    11, WHITE, spacing=1.15)
footer(s, "checkpoints/  ·  last.pt @ epoch 7  ·  best_val 3.47")

prs.save("NJEM_detailed.pptx")
print(f"wrote NJEM_detailed.pptx ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
