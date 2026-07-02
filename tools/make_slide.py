from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ── Modern light palette ────────────────────────────────
INK    = RGBColor(0x0F, 0x17, 0x2A)
SLATE  = RGBColor(0x33, 0x41, 0x55)
MUTED  = RGBColor(0x64, 0x74, 0x8B)
BORDER = RGBColor(0xE2, 0xE8, 0xF0)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
BG_A   = RGBColor(0xFB, 0xFC, 0xFF)
BG_B   = RGBColor(0xEA, 0xF0, 0xFD)
CARD_A = RGBColor(0xFF, 0xFF, 0xFF)
CARD_B = RGBColor(0xF3, 0xF6, 0xFD)
ACCENTS = [
    (RGBColor(0x63,0x66,0xF1), RGBColor(0x81,0x8C,0xF8), RGBColor(0x4F,0x46,0xE5)),  # indigo
    (RGBColor(0x8B,0x5C,0xF6), RGBColor(0xA7,0x8B,0xFA), RGBColor(0x7C,0x3A,0xED)),  # violet
    (RGBColor(0x06,0xB6,0xD4), RGBColor(0x22,0xD3,0xEE), RGBColor(0x08,0x91,0xB2)),  # cyan
    (RGBColor(0x14,0xB8,0xA6), RGBColor(0x2D,0xD4,0xBF), RGBColor(0x0D,0x94,0x88)),  # teal
    (RGBColor(0xF5,0x9E,0x0B), RGBColor(0xFB,0xBF,0x24), RGBColor(0xD9,0x77,0x06)),  # amber
    (RGBColor(0xEC,0x48,0x99), RGBColor(0xF4,0x72,0xB6), RGBColor(0xDB,0x27,0x77)),  # pink
]

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def grad(fill, c1, c2, angle=90):
    fill.gradient()
    s = fill.gradient_stops
    s[0].position = 0.0; s[0].color.rgb = c1
    s[1].position = 1.0; s[1].color.rgb = c2
    try: fill.gradient_angle = angle
    except Exception: pass

def new_slide():
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    grad(sl.background.fill, BG_A, BG_B, angle=115)
    return sl

def txt(sl, x, y, w, h, lines, size, color, bold=False, italic=False,
        align=PP_ALIGN.LEFT, font="Calibri", anchor=MSO_ANCHOR.TOP, spacing=1.0, cspace=None):
    tb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
    if isinstance(lines, str): lines = [lines]
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = spacing
        r = p.add_run(); r.text = ln
        f = r.font; f.size = Pt(size); f.bold = bold; f.italic = italic
        f.name = font; f.color.rgb = color
        if cspace is not None:
            r._r.get_or_add_rPr().set("spc", str(cspace))
    return tb

def soft_shadow(sp):
    spPr = sp._element.spPr
    eff = spPr.makeelement(qn('a:effectLst'), {})
    sh = spPr.makeelement(qn('a:outerShdw'),
                          {'blurRad':'90000','dist':'40000','dir':'5400000','rotWithShape':'0'})
    clr = spPr.makeelement(qn('a:srgbClr'), {'val':'9FB0CC'})
    clr.append(spPr.makeelement(qn('a:alpha'), {'val':'38000'}))
    sh.append(clr); eff.append(sh); spPr.append(eff)

def card(sl, x, y, w, h, c1=CARD_A, c2=CARD_B, angle=90, radius=0.10, border=BORDER, shadow=True):
    sp = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    grad(sp.fill, c1, c2, angle)
    if border is not None: sp.line.color.rgb = border; sp.line.width = Pt(1)
    else: sp.line.fill.background()
    try: sp.adjustments[0] = radius
    except Exception: pass
    sp.shadow.inherit = False
    if shadow: soft_shadow(sp)
    return sp

def circle(sl, x, y, d, c1, c2):
    sp = sl.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    grad(sp.fill, c1, c2, angle=45); sp.line.fill.background(); sp.shadow.inherit = False
    soft_shadow(sp); return sp

# ════════════════════════════════════════════════════════
# SLIDE 1 — Overview
# ════════════════════════════════════════════════════════
s1 = new_slide()
txt(s1, 0.6, 0.42, 12.1, 0.7, "Decoding Speech from Brain Activity", 34, INK, bold=True, font="Cambria")
txt(s1, 0.62, 1.14, 12.1, 0.45,
    "Mapping intracortical neural signals to Whisper audio embeddings, then to text",
    16, ACCENTS[0][2], italic=True)

stages = [
    ("Neural Activity", "(T x 512)", ["Intracortical t15", "recordings, 20 ms bins"]),
    ("ConvBiGRU Model", "12.3M params", ["Conv + 8-layer BiGRU", "+ cross-attention"]),
    ("Whisper Embeddings", "(1500 x 384)", ["Predicted encoder", "feature sequence"]),
    ("Text Transcript", "frozen decoder", ["Whisper tiny.en decodes", "embeddings -> words"]),
]
cardW, cardH, gap = 2.75, 2.35, 0.42
x0 = (13.333 - (4*cardW + 3*gap)) / 2
yC = 2.1
for i, (t, sp_, d) in enumerate(stages):
    a1, a2, adk = ACCENTS[i]
    x = x0 + i*(cardW+gap)
    card(s1, x, yC, cardW, cardH)
    cd = 0.66
    circle(s1, x+cardW/2-cd/2, yC+0.26, cd, a1, a2)
    txt(s1, x+cardW/2-cd/2, yC+0.26, cd, cd, str(i+1), 22, WHITE, bold=True,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font="Cambria")
    txt(s1, x+0.1, yC+1.05, cardW-0.2, 0.55, t, 15, INK, bold=True, align=PP_ALIGN.CENTER, font="Cambria")
    txt(s1, x, yC+1.58, cardW, 0.3, sp_, 12, adk, bold=True, align=PP_ALIGN.CENTER, font="Courier New")
    txt(s1, x+0.12, yC+1.92, cardW-0.24, 0.7, d, 11, MUTED, align=PP_ALIGN.CENTER, spacing=0.95)
    if i < 3:
        txt(s1, x+cardW-0.02, yC+0.9, gap+0.06, 0.6, "→", 24, adk, bold=True,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

facts = [
    ("Target", ["Whisper encoder embeddings,", "length-masked to valid frames"]),
    ("Loss", ["SmoothL1 + (1 - cosine)", "+ decoder-in-the-loop CE"]),
    ("Eval", ["WER by decoding predicted", "embeddings through Whisper"]),
]
fW, fGap = 3.9, 0.3
fx0 = (13.333 - (3*fW + 2*fGap)) / 2
fy = 5.0
for i, (k, v) in enumerate(facts):
    a1, a2, adk = ACCENTS[i]
    x = fx0 + i*(fW+fGap)
    card(s1, x, fy, fW, 1.5)
    circle(s1, x+0.28, fy+0.3, 0.34, a1, a2)
    txt(s1, x+0.78, fy+0.24, fW-0.95, 0.4, k, 15, INK, bold=True, font="Cambria")
    txt(s1, x+0.3, fy+0.72, fW-0.55, 0.7, v, 12, MUTED, spacing=0.95)

txt(s1, 0.6, 6.78, 12.1, 0.4, "t15 Brain-to-Text   |   NJEM — Neural Joint Embedding Model",
    12, MUTED, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════
# SLIDE 2 — Architecture & Computation
# ════════════════════════════════════════════════════════
s2 = new_slide()
txt(s2, 0.6, 0.42, 12.1, 0.7, "Model Architecture & Computation", 34, INK, bold=True, font="Cambria")
txt(s2, 0.62, 1.14, 12.1, 0.45,
    "ConvBiGRU forward pass (neural -> embeddings -> text) and the training objective",
    16, ACCENTS[0][2], italic=True)

# ── Left column: forward pass ───────────────────────────
txt(s2, 0.62, 1.75, 7.0, 0.35, "FORWARD PASS", 13, ACCENTS[0][2], bold=True, cspace=180)
blocks = [
    ("Input  —  neural features", "x : (B, T, 512)   z-scored, 20 ms bins"),
    ("Conv front-end", "2 x Conv1d(k=5, s=2) + GELU + Drop  ->  (B, T/4, 256)"),
    ("BiGRU (packed sequences)", "8 layers, hidden 256, bidirectional  ->  (B, T', 512)"),
    ("Cross-attention", "N learnable queries, 4 heads, key-pad mask  ->  (B, N, 512)"),
    ("LayerNorm + MLP head", "Linear 512->512, GELU, Linear 512->384  ->  (B, N, 384)"),
    ("Frozen Whisper decoder", "predicted embeddings  ->  text tokens (inference)"),
]
lx, lw = 0.62, 6.9
by = 2.18; bh = 0.7; bgap = 0.075
for i, (t, sub) in enumerate(blocks):
    a1, a2, adk = ACCENTS[i]
    y = by + i*(bh+bgap)
    card(s2, lx, y, lw, bh)
    cd = 0.42
    circle(s2, lx+0.2, y+bh/2-cd/2, cd, a1, a2)
    txt(s2, lx+0.2, y+bh/2-cd/2, cd, cd, str(i+1), 14, WHITE, bold=True,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font="Cambria")
    txt(s2, lx+0.8, y+0.1, lw-0.95, 0.32, t, 13.5, INK, bold=True, font="Cambria")
    txt(s2, lx+0.8, y+0.4, lw-0.95, 0.28, sub, 10.5, MUTED, font="Courier New")

# ── Right column: training objective ────────────────────
rx, rw = 7.85, 4.85
txt(s2, rx, 1.75, rw, 0.35, "TRAINING OBJECTIVE", 13, ACCENTS[3][2], bold=True, cspace=180)

# valid-frames note
card(s2, rx, 2.18, rw, 0.72)
txt(s2, rx+0.25, 2.3, rw-0.5, 0.28, "Frame masking", 13, INK, bold=True, font="Cambria")
txt(s2, rx+0.25, 2.58, rw-0.5, 0.28, "vf = ceil(audio_length / 320)  (content frames)",
    10.5, MUTED, font="Courier New")

loss_cards = [
    (2, "Embedding loss  (masked)",
        ["SmoothL1(pred, tgt) + (1 - cos(pred, tgt))", "averaged over valid frames only"]),
    (4, "Decoder-in-the-loop  CE",
        ["teacher-force GT tokens through the", "FROZEN Whisper decoder on predictions"]),
    (5, "Total loss",
        ["L = L_emb + λ · L_CE", "λ = dec_loss_weight = 0.5"]),
]
ry = 3.08; rh = 1.15; rgap = 0.14
for j, (ai, t, body) in enumerate(loss_cards):
    a1, a2, adk = ACCENTS[ai]
    y = ry + j*(rh+rgap)
    card(s2, rx, y, rw, rh)
    circle(s2, rx+0.28, y+0.3, 0.34, a1, a2)
    txt(s2, rx+0.78, y+0.22, rw-0.95, 0.35, t, 14, INK, bold=True, font="Cambria")
    fnt = "Courier New" if j == 2 else "Calibri"
    txt(s2, rx+0.3, y+0.62, rw-0.55, 0.5, body, 11, adk if j == 2 else MUTED,
        spacing=0.98, font=fnt)

txt(s2, 0.6, 6.95, 12.1, 0.35,
    "Whisper tiny.en:  n_audio_ctx = 1500 frames  |  n_audio_state = 384  |  encoder frozen throughout",
    11, MUTED, align=PP_ALIGN.CENTER)

prs.save("NJEM_overview.pptx")
print("wrote NJEM_overview.pptx (2 slides)")
